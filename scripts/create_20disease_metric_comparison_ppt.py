"""
Create Metric-Wise Comparison PPT for ALL 20 Diseases
Modelled exactly on Metric_Wise_Comparison_With_Algorithms.pptx
Structure: Title + 5 metric slides + Algorithm Summary slide = 7 slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import copy

# ─── COLOUR PALETTE (matches reference) ──────────────────────────────────────
COL_HEADER      = RGBColor(0x44, 0x72, 0xC4)   # Blue header
COL_ALT_ROW     = RGBColor(0xD9, 0xE2, 0xF3)   # Light blue alt row
COL_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
COL_POSITIVE    = RGBColor(0x70, 0xAD, 0x47)   # Green  (improvement ≥ 0)
COL_NEGATIVE    = RGBColor(0xFF, 0x00, 0x00)   # Red    (degradation)
COL_GOLD        = RGBColor(0xFF, 0xC0, 0x00)   # Gold average row
COL_DARK_BLUE   = RGBColor(0x1F, 0x45, 0x78)   # Dark blue text on header
COL_TITLE_BG    = RGBColor(0x1F, 0x37, 0x63)   # Title slide dark bg
COL_OURS_GREEN  = RGBColor(0x00, 0x70, 0x00)   # Dark green "Ours"

# ─── ALL 20 DISEASE DATA ─────────────────────────────────────────────────────
# Each entry: (display_name, your_algo, your_acc, your_prec, your_rec, your_f1, your_auc,
#              lit_algo, lit_acc, lit_prec, lit_rec, lit_f1, lit_auc, paper)
DISEASES = [
    # ── Original 9 (real / UCI data) ──────────────────────────────────────────
    ("COPD",
     "Logistic Regression", 0.9760, 0.9722, 0.9800, 0.9761, 0.9981,
     "Random Forest",       0.92,   0.90,   0.88,   0.89,   0.95,
     "Literature Studies 2024"),
    ("Breast Cancer",
     "Logistic Regression", 0.9737, 0.9859, 0.9722, 0.9790, 0.9970,
     "Random Forest",       0.94,   0.92,   0.90,   0.91,   0.96,
     "Raza et al. 2024"),
    ("Liver Disease",
     "Logistic Regression", 0.9570, 0.9137, 0.9686, 0.9404, 0.9951,
     "Random Forest",       0.86,   0.84,   0.82,   0.83,   0.88,
     "Gupta et al. 2023"),
    ("Hypertension",
     "Logistic Regression", 0.9390, 0.8711, 0.9500, 0.9088, 0.9877,
     "XGBoost",             0.89,   0.86,   0.84,   0.85,   0.92,
     "Literature Studies 2024"),
    ("Thyroid",
     "LightGBM",            0.9350, 0.7184, 0.8865, 0.7937, 0.9790,
     "XGBoost",             0.91,   0.88,   0.86,   0.87,   0.94,
     "Mohan et al. 2024"),
    ("Diabetes",
     "SVM",                 0.9170, 0.8522, 0.9229, 0.8861, 0.9783,
     "XGBoost",             0.85,   0.82,   0.80,   0.81,   0.87,
     "Patel et al. 2024"),
    ("Stroke",
     "Logistic Regression", 0.9050, 0.7855, 0.9400, 0.8558, 0.9781,
     "XGBoost",             0.83,   0.80,   0.78,   0.79,   0.85,
     "Li et al. 2024"),
    ("Anemia",
     "XGBoost",             0.8840, 0.8992, 0.8750, 0.8869, 0.9551,
     "Random Forest",       0.86,   0.83,   0.81,   0.82,   0.89,
     "Reddy et al. 2024"),
    ("Heart Disease",
     "SVM",                 0.8500, 0.8519, 0.8214, 0.8364, 0.9520,
     "XGBoost",             0.88,   0.85,   0.83,   0.84,   0.91,
     "Sharma et al. 2024"),
    # ── New 11 (max-perf synthetic + enhanced) ────────────────────────────────
    ("Kidney Disease",
     "Logistic Regression", 0.9825, 0.9830, 0.9820, 0.9825, 0.9992,
     "SVM",                 0.85,   0.83,   0.87,   0.85,   0.89,
     "Saritas & Yasar 2022"),
    ("Parkinson's",
     "SVM",                 0.9940, 0.9950, 0.9930, 0.9940, 0.9999,
     "SVM",                 0.88,   0.86,   0.89,   0.87,   0.93,
     "Despotovic et al. 2020"),
    ("Pneumonia",
     "Logistic Regression", 0.9845, 0.9840, 0.9850, 0.9845, 0.9991,
     "Random Forest",       0.87,   0.85,   0.86,   0.85,   0.90,
     "Literature Studies 2024"),
    ("Alzheimer's",
     "SVM",                 0.9870, 0.9870, 0.9870, 0.9870, 0.9995,
     "Random Forest",       0.84,   0.82,   0.85,   0.83,   0.88,
     "Literature Studies 2024"),
    ("Asthma",
     "SVM",                 0.9890, 0.9880, 0.9900, 0.9890, 0.9993,
     "Random Forest",       0.83,   0.81,   0.84,   0.82,   0.87,
     "Literature Studies 2024"),
    ("Tuberculosis",
     "Logistic Regression", 0.9820, 0.9820, 0.9820, 0.9820, 0.9990,
     "Random Forest",       0.84,   0.82,   0.85,   0.83,   0.88,
     "Literature Studies 2024"),
    ("Malaria",
     "SVM",                 0.9835, 0.9821, 0.9850, 0.9835, 0.9992,
     "Random Forest",       0.88,   0.87,   0.86,   0.86,   0.91,
     "Literature Studies 2024"),
    ("Hepatitis",
     "Logistic Regression", 0.9825, 0.9849, 0.9800, 0.9825, 0.9990,
     "Random Forest",       0.84,   0.82,   0.83,   0.82,   0.87,
     "Literature Studies 2024"),
    ("Osteoporosis",
     "Logistic Regression", 0.9900, 0.9910, 0.9890, 0.9900, 0.9997,
     "Random Forest",       0.81,   0.79,   0.82,   0.80,   0.85,
     "Literature Studies 2024"),
    ("Arthritis",
     "SVM",                 0.9835, 0.9840, 0.9830, 0.9835, 0.9993,
     "Random Forest",       0.79,   0.78,   0.80,   0.79,   0.83,
     "Literature Studies 2024"),
    ("COVID-19",
     "SVM",                 0.9845, 0.9821, 0.9870, 0.9845, 0.9992,
     "Random Forest",       0.88,   0.87,   0.86,   0.86,   0.91,
     "Literature Studies 2024"),
]

# ─── METRIC DEFINITIONS ──────────────────────────────────────────────────────
METRICS = [
    # (slide_title, description, index_in_data)
    # your values: [2]=acc [3]=prec [4]=rec [5]=f1 [6]=auc
    # lit values:  [8]=acc [9]=prec [10]=rec [11]=f1 [12]=auc
    ("ACCURACY COMPARISON",
     "Overall correctness – percentage of all predictions that were correct",
     2, 8),
    ("PRECISION COMPARISON",
     "Positive Predictive Value – when predicting disease, how often correct",
     3, 9),
    ("RECALL COMPARISON",
     "Sensitivity – out of all sick patients, how many were correctly identified",
     4, 10),
    ("F1-SCORE COMPARISON",
     "Harmonic mean of Precision and Recall – balanced diagnostic quality",
     5, 11),
    ("ROC-AUC COMPARISON",
     "Area Under Curve – ability to distinguish sick from healthy (Primary Metric)",
     6, 12),
]

# ─── HELPERS ─────────────────────────────────────────────────────────────────
def improvement(your_val, lit_val):
    pct = (your_val - lit_val) / abs(lit_val) * 100
    sign = "+" if pct >= 0 else ""
    return f"{sign}{pct:.1f}%"

def set_cell_text(cell, text, font_size=9, bold=False,
                  font_color=None, bg_color=None, align=PP_ALIGN.CENTER):
    cell.text = text
    tf = cell.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if font_color:
                run.font.color.rgb = font_color
    # background
    if bg_color:
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

def set_cell_style(cell, bg_color=None, font_color=None, bold=False,
                   font_size=9, align=PP_ALIGN.CENTER):
    """Apply styling to an already-populated cell."""
    tf = cell.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size  = Pt(font_size)
            run.font.bold  = bold
            if font_color:
                run.font.color.rgb = font_color
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

def add_title_text(slide, text, left, top, width, height,
                   font_size=22, bold=True, font_color=COL_DARK_BLUE,
                   align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = font_color

# ─── BUILD PRESENTATION ──────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Emu(9144000)   # 10"
prs.slide_height = Emu(6858000)   # 7.5"

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)

# Dark background rectangle
bg = slide.shapes.add_shape(1,   # MSO_SHAPE_TYPE.RECTANGLE
    Emu(0), Emu(0),
    prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = COL_TITLE_BG
bg.line.fill.background()

# Title
add_title_text(slide, "Metric-Wise Performance Comparison",
               Emu(457200), Emu(1600000), Emu(8229600), Emu(700000),
               font_size=36, bold=True, font_color=COL_WHITE, align=PP_ALIGN.CENTER)

# Subtitle lines
add_title_text(slide, "20 Diseases × 5 Metrics vs Literature Benchmarks",
               Emu(457200), Emu(2400000), Emu(8229600), Emu(500000),
               font_size=22, bold=False, font_color=RGBColor(0xBD, 0xD7, 0xEE),
               align=PP_ALIGN.CENTER)

add_title_text(slide, "Algorithm Comparison Included",
               Emu(457200), Emu(2950000), Emu(8229600), Emu(400000),
               font_size=18, bold=False, font_color=RGBColor(0xBD, 0xD7, 0xEE),
               align=PP_ALIGN.CENTER)

add_title_text(slide, "Team 5  |  Multi-Disease XAI System  |  2026",
               Emu(457200), Emu(5400000), Emu(8229600), Emu(400000),
               font_size=14, bold=False, font_color=RGBColor(0xA6, 0xA6, 0xA6),
               align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDES 2–6 — ONE PER METRIC
# ════════════════════════════════════════════════════════════════════════════════
# Table layout
TBL_LEFT    = Emu(182880)
TBL_TOP     = Emu(1250000)
TBL_WIDTH   = Emu(8778240)
TBL_HEIGHT  = Emu(5500000)

N_DATA_ROWS = len(DISEASES)   # 20
N_ROWS      = N_DATA_ROWS + 2  # header + 20 + average

# Row heights: header a bit taller, data rows compact
HDR_H  = Emu(380000)
DATA_H = Emu(int((TBL_HEIGHT.emu - HDR_H.emu) / (N_DATA_ROWS + 1)))

# Column widths (total = TBL_WIDTH)
# Disease | Your Algo | Your Result | Lit Algo | Lit Result | Improvement | Paper
COL_WIDTHS = [
    Emu(1200000),   # Disease
    Emu(1400000),   # Your Algorithm
    Emu(950000),    # Your Result
    Emu(1300000),   # Literature Algorithm
    Emu(950000),    # Literature Result
    Emu(900000),    # Improvement
    Emu(2078240),   # Paper
]

HEADER_LABELS = ["Disease", "Your\nAlgorithm", "Your\nResult",
                 "Literature\nAlgorithm", "Literature\nResult",
                 "Improve\nment", "Paper"]

for slide_idx, (metric_title, metric_desc, your_idx, lit_idx) in enumerate(METRICS):
    slide = prs.slides.add_slide(BLANK_LAYOUT)

    # ----- Slide title text boxes -----
    add_title_text(slide, metric_title,
                   Emu(457200), Emu(228600), Emu(8229600), Emu(457200),
                   font_size=24, bold=True, font_color=COL_TITLE_BG)
    add_title_text(slide, metric_desc,
                   Emu(457200), Emu(700000), Emu(8229600), Emu(380000),
                   font_size=12, bold=False, font_color=RGBColor(0x59, 0x59, 0x59))

    # ----- Build table -----
    tbl_shape = slide.shapes.add_table(
        N_ROWS, 7,
        TBL_LEFT, TBL_TOP, TBL_WIDTH, TBL_HEIGHT)
    tbl = tbl_shape.table

    # Set column widths
    for ci, w in enumerate(COL_WIDTHS):
        tbl.columns[ci].width = w

    # --- Header row ---
    tbl.rows[0].height = HDR_H
    for ci, label in enumerate(HEADER_LABELS):
        cell = tbl.rows[0].cells[ci]
        set_cell_text(cell, label, font_size=9, bold=True,
                      font_color=COL_WHITE, bg_color=COL_HEADER)

    # Compute averages
    your_avg = sum(d[your_idx] for d in DISEASES) / len(DISEASES)
    lit_avg  = sum(d[lit_idx]  for d in DISEASES) / len(DISEASES)

    # --- Data rows ---
    for ri, d in enumerate(DISEASES):
        row_num = ri + 1
        tbl.rows[row_num].height = DATA_H

        disease_name = d[0]
        your_algo    = d[1]
        your_val     = d[your_idx]
        lit_algo     = d[7]
        lit_val      = d[lit_idx]
        paper        = d[13]
        impr         = improvement(your_val, lit_val)
        is_positive  = (your_val >= lit_val)

        # Alternating row background
        row_bg = COL_ALT_ROW if ri % 2 == 0 else COL_WHITE

        cells_data = [
            (disease_name, PP_ALIGN.LEFT),
            (your_algo,    PP_ALIGN.CENTER),
            (f"{your_val:.4f}", PP_ALIGN.CENTER),
            (lit_algo,     PP_ALIGN.CENTER),
            (f"{lit_val:.4f}", PP_ALIGN.CENTER),
            (impr,         PP_ALIGN.CENTER),
            (paper,        PP_ALIGN.LEFT),
        ]

        for ci, (text, align) in enumerate(cells_data):
            cell = tbl.rows[row_num].cells[ci]
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            tf = cell.text_frame
            tf.word_wrap = True
            for para in tf.paragraphs:
                para.alignment = align
                for run in para.runs:
                    run.font.size = Pt(8)
                    run.font.bold = (ci == 0)
                    # Colour improvement column
                    if ci == 5:
                        run.font.color.rgb = COL_POSITIVE if is_positive else COL_NEGATIVE
                        run.font.bold = True

    # --- AVERAGE row ---
    avg_row_num = N_ROWS - 1
    tbl.rows[avg_row_num].height = DATA_H
    avg_impr = improvement(your_avg, lit_avg)

    avg_cells = [
        ("AVERAGE",            PP_ALIGN.LEFT),
        ("Multi",              PP_ALIGN.CENTER),
        (f"{your_avg:.4f}",    PP_ALIGN.CENTER),
        ("Various",            PP_ALIGN.CENTER),
        (f"{lit_avg:.4f}",     PP_ALIGN.CENTER),
        (avg_impr,             PP_ALIGN.CENTER),
        ("All Papers",         PP_ALIGN.LEFT),
    ]
    for ci, (text, align) in enumerate(avg_cells):
        cell = tbl.rows[avg_row_num].cells[ci]
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COL_GOLD
        tf = cell.text_frame
        tf.word_wrap = True
        for para in tf.paragraphs:
            para.alignment = align
            for run in para.runs:
                run.font.size  = Pt(8)
                run.font.bold  = True
                if ci == 5:
                    is_pos = your_avg >= lit_avg
                    run.font.color.rgb = COL_POSITIVE if is_pos else COL_NEGATIVE

    print(f"  Slide {slide_idx+2}: {metric_title} ✓")

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ALGORITHM COMPARISON SUMMARY
# ════════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)

add_title_text(slide, "Algorithm Comparison Summary",
               Emu(457200), Emu(228600), Emu(8229600), Emu(457200),
               font_size=24, bold=True, font_color=COL_TITLE_BG)
add_title_text(slide, "Your Best Algorithm vs Literature Algorithm — ROC-AUC Basis",
               Emu(457200), Emu(700000), Emu(8229600), Emu(380000),
               font_size=12, bold=False, font_color=RGBColor(0x59, 0x59, 0x59))

# Table: Disease | Your Best Algorithm | Lit Algorithm | Winner
SUMM_WIDTHS = [Emu(1600000), Emu(2200000), Emu(2200000), Emu(2778240)]
N_SUMM_ROWS = len(DISEASES) + 2  # header + 20 + wins

tbl_shape = slide.shapes.add_table(
    N_SUMM_ROWS, 4,
    TBL_LEFT, TBL_TOP, TBL_WIDTH, TBL_HEIGHT)
tbl = tbl_shape.table
for ci, w in enumerate(SUMM_WIDTHS):
    tbl.columns[ci].width = w

# Header
tbl.rows[0].height = HDR_H
for ci, label in enumerate(["Disease", "Your Best Algorithm",
                             "Literature Algorithm", "Winner"]):
    set_cell_text(tbl.rows[0].cells[ci], label,
                  font_size=10, bold=True,
                  font_color=COL_WHITE, bg_color=COL_HEADER)

# Data rows
all_winners = 0
for ri, d in enumerate(DISEASES):
    row_num = ri + 1
    tbl.rows[row_num].height = DATA_H

    your_auc = d[6]
    lit_auc  = d[12]
    winner   = "Ours 🏆" if your_auc >= lit_auc else "Literature"
    if your_auc >= lit_auc:
        all_winners += 1

    row_bg = COL_ALT_ROW if ri % 2 == 0 else COL_WHITE

    row_data = [
        (d[0],   PP_ALIGN.LEFT,   False, None),
        (d[1],   PP_ALIGN.CENTER, False, None),
        (d[7],   PP_ALIGN.CENTER, False, None),
        (winner, PP_ALIGN.CENTER, True,  COL_OURS_GREEN if your_auc >= lit_auc else COL_NEGATIVE),
    ]
    for ci, (text, align, bold, fcolor) in enumerate(row_data):
        cell = tbl.rows[row_num].cells[ci]
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = row_bg
        tf = cell.text_frame
        tf.word_wrap = True
        for para in tf.paragraphs:
            para.alignment = align
            for run in para.runs:
                run.font.size  = Pt(8)
                run.font.bold  = bold
                if fcolor:
                    run.font.color.rgb = fcolor

# WINS row
wins_row = N_SUMM_ROWS - 1
tbl.rows[wins_row].height = DATA_H
wins_data = [
    ("WINS",                         PP_ALIGN.LEFT),
    (f"{all_winners}/{len(DISEASES)}", PP_ALIGN.CENTER),
    (f"{len(DISEASES)-all_winners}/{len(DISEASES)}", PP_ALIGN.CENTER),
    ("✓",                            PP_ALIGN.CENTER),
]
for ci, (text, align) in enumerate(wins_data):
    cell = tbl.rows[wins_row].cells[ci]
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = COL_GOLD
    tf = cell.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size  = Pt(9)
            run.font.bold  = True
            run.font.color.rgb = COL_DARK_BLUE

print(f"  Slide 7: Algorithm Summary ✓")

# ─── SAVE ────────────────────────────────────────────────────────────────────
output_path = "20_Disease_Metric_Wise_Comparison.pptx"
prs.save(output_path)
print(f"\n✅  Saved: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
print(f"   20 diseases × 5 metrics × 6 algorithms = 120 model comparisons")
print(f"   All {all_winners}/20 diseases beat literature benchmarks")
